"""Webhook handlers for Telegram and WhatsApp."""
from fastapi import APIRouter, Request, HTTPException, BackgroundTasks
from sqlalchemy import text
from app.config import settings
from app.database import async_session_factory
from app.services.agent_manager import hermes_profile_chat_with_fallback
import os
import re
from app.services.profile_init import init_user_profile
import httpx, json, asyncio
import sys, subprocess, time
from datetime import datetime, timedelta, timezone
from zoneinfo import ZoneInfo
from pathlib import Path
from app.logging_config import get_logger
from app.services.audit_logger import audit_logger, AuditLogger

logger = get_logger(__name__)

router = APIRouter()

WELCOME = (
    "🎉 **Welcome! Your AI agent is ready.**\n\n"
    "Here's what I can do:\n"
    "📝 *\"Save a note about...\"* — I'll store it in your vault\n"
    "📖 *\"What notes do I have?\"* — I'll read your vault\n"
    "🌐 *\"Search for...\"* — I'll find and summarize\n"
    "💾 I remember our conversations\n\n"
    "👉 **Try saying:** *\"Save a note about my meeting today\"* or *\"Search for the latest AI news\"*"
)

def _safe_chat_id(chat_id: str) -> int:
    """Safely convert chat_id to int with validation.

    SECURITY FIX: Prevents uncaught ValueError from invalid chat_id.
    """
    try:
        # Strip whitespace and validate it's a valid integer (can be negative for groups)
        return int(str(chat_id).strip())
    except (ValueError, TypeError):
        raise ValueError(f"Invalid chat_id format: {chat_id}")


async def typing_indicator(chat_id: str, stop_event: asyncio.Event):
    """Keep typing indicator alive every 4 seconds until stop_event is set."""
    try:
        safe_id = _safe_chat_id(chat_id)
    except ValueError:
        return  # Invalid chat_id, silently stop typing indicator

    while not stop_event.is_set():
        try:
            async with httpx.AsyncClient() as c:
                await c.post(
                    f"https://api.telegram.org/bot{settings.telegram_bot_token}/sendChatAction",
                    json={"chat_id": safe_id, "action": "typing"},
                    timeout=5,
                )
        except: pass
        try:
            await asyncio.wait_for(stop_event.wait(), timeout=4)
        except asyncio.TimeoutError:
            continue

async def send_tg(chat_id: str, text: str):
    if not text:
        return
    try:
        safe_id = _safe_chat_id(chat_id)
    except ValueError as e:
        logger.warning("send_tg_invalid_chat_id", chat_id=chat_id, error=str(e))
        return

    if len(text) > 4000:  # Telegram hard limit; truncate to avoid sendMessage 400
        text = text[:3975] + "\n…(truncated)"
    async with httpx.AsyncClient() as c:
        await c.post(
            f"https://api.telegram.org/bot{settings.telegram_bot_token}/sendMessage",
            json={"chat_id": safe_id, "text": text},  # plain text: robust to markdown/emoji
        )


async def send_tg_file(chat_id: str, path: str, caption: str = ""):
    """Send a file (e.g. .pptx) to Telegram as a document attachment."""
    try:
        safe_id = _safe_chat_id(chat_id)
    except ValueError as e:
        logger.warning("send_tg_file_invalid_chat_id", chat_id=chat_id, error=str(e))
        return

    async with httpx.AsyncClient(timeout=60) as c:
        with open(path, "rb") as f:
            await c.post(
                f"https://api.telegram.org/bot{settings.telegram_bot_token}/sendDocument",
                data={"chat_id": safe_id, "caption": caption},
                files={"document": (os.path.basename(path), f,
                                    "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
            )


def _deck_text(block) -> str:
    """Flatten a slide block into plain text lines for the pptx body."""
    if isinstance(block, str):
        return block
    if isinstance(block, list):
        return "\n".join(str(x) for x in block)
    if isinstance(block, dict):
        # common keys: bullets/content/text/items/notes/subtitle
        for k in ("bullets", "content", "items", "text", "points"):
            if block.get(k):
                return _deck_text(block[k])
        return _deck_text(list(block.values())[0] if block else "")
    return ""


def build_pptx(deck: dict) -> str:
    """Convert a deck spec (title/slides) into a .pptx file, return the path.
    Uses the blank layout + explicit textboxes so it tolerates any slide layout
    (title/section/content) without placeholder-index failures."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)  # 16:9
    blank = prs.slide_layouts[6]  # Blank
    for s in deck.get("slides", []):
        if not isinstance(s, dict):
            continue
        slide = prs.slides.add_slide(blank)
        title = str(s.get("title", "")).strip()
        body_txt = _deck_text(
            s.get("content") or s.get("bullets") or s.get("subtitle")
            or s.get("points") or s.get("body") or s.get("notes") or ""
        )
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.3))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = title
        run.font.size = Pt(32)
        run.font.bold = True
        if body_txt:
            tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.7), Inches(5.2))
            tf = tb2.text_frame
            tf.word_wrap = True
            tf.text = body_txt
            for par in tf.paragraphs:
                for r in par.runs:
                    if r.font.size is None:
                        r.font.size = Pt(18)
    out = f"/tmp/{int(__import__('time').time())}_deck.pptx"
    prs.save(out)
    return out


PPTX_CREATE = "/usr/local/lib/hermes-agent/skills/productivity/powerpoint/scripts/pptx_create.py"


def _render_deck(data: dict):
    """Render a deck JSON spec to a real .pptx using the Hermes 'powerpoint' skill
    (pptx_create.py), which supports backgrounds/footers/slide numbers/notes."""
    if not os.path.exists(PPTX_CREATE):
        return None
    fin = f"/tmp/deck_in_{int(__import__('time').time())}.json"
    out = f"/tmp/deck_out_{int(__import__('time').time())}.pptx"
    try:
        with open(fin, "w") as f:
            json.dump(data, f)
        r = subprocess.run([sys.executable, PPTX_CREATE, fin, out],
                           capture_output=True, timeout=120)
        return out if (r.returncode == 0 and os.path.exists(out)) else None
    except Exception:
        return None
    finally:
        if os.path.exists(fin):
            os.remove(fin)


def _parse_deck(arg: str):
    """Extract a deck dict (has a non-empty 'slides' list) from raw or diff-wrapped agent output."""
    if not arg:
        return None
    data = None
    try:
        data = json.loads(arg)
    except Exception:
        data = None
    if not isinstance(data, dict):
        clean = "\n".join(
            (l[1:] if l.startswith("+") else l)
            for l in arg.splitlines()
            if l.strip() and not l.startswith(("@", "a/", "b/", "diff ", "index ", "\\"))
        )
        i, j = clean.find("{"), clean.rfind("}")
        if i != -1 and j > i:
            try:
                data = json.loads(clean[i:j + 1])
            except Exception:
                data = None
    if isinstance(data, dict) and isinstance(data.get("slides"), list) and data["slides"]:
        return data
    return None


def try_build_deck(arg: str):
    """If arg is a deck-JSON, render it to a .pptx via the powerpoint skill and return the path."""
    data = _parse_deck(arg)
    if data is None:
        return None
    try:
        return _render_deck(data)
    except Exception:
        return None


def _deck_summary(data: dict) -> str:
    """Plain-prose outline of a deck (never raw JSON) for the fallback chat reply."""
    lines = [f"*{data.get('title') or 'Your slideshow'}*", ""]
    for i, s in enumerate(data.get("slides", []), 1):
        st = (s.get("title") if isinstance(s, dict) else None) or (s.get("subtitle") if isinstance(s, dict) else None) or f"Slide {i}"
        lines.append(f"{i}. {str(st)[:120]}")
    lines.append("")
    lines.append("I prepared these slides. I couldn't attach the file just now — ask me again and I'll make sure you get the .pptx.")
    return "\n".join(lines)


async def deliver_reply(chat_id: str, resp: str):
    """Send the agent reply. If it's a deck, deliver a real .pptx document; NEVER send deck JSON as text."""
    if not resp:
        return
    deck_path = try_build_deck(resp)
    if deck_path:
        try:
            await send_tg(chat_id, "📊 Here's your slideshow.")
            await send_tg_file(chat_id, deck_path, caption="Research slides (generated by your agent)")
        finally:
            try:
                os.remove(deck_path)
            except Exception:
                pass
        return
    data = _parse_deck(resp)
    if data is not None:
        await send_tg(chat_id, _deck_summary(data))
        return
    if '"slides"' in resp or '"slide_size"' in resp:  # looks like a deck we failed to parse — never dump raw JSON
        await send_tg(chat_id, "I prepared a slideshow for you, but I ran into trouble attaching the file. Ask me again (e.g. \"make that a slideshow again\") and I'll send you the .pptx.")
        return
    if _looks_like_raw_code(resp):
        await send_tg(chat_id, "I worked that out as a script for you. If you'd like to see it, just say \"show me the code\" and I'll paste it.")
        return
    await send_tg(chat_id, resp)

def _looks_like_raw_code(text: str) -> bool:
    """Detect a reply that is primarily raw generated code / a file diff (not a
    user-facing answer). Used to stop python/script dumps from leaking to the user."""
    s = text.lstrip()
    if re.search(r"^```(?:python|py)\s*$", s, flags=re.M):
        return True
    if re.search(r"^(?:def |import |from |class |#!)", s) and len(s) > 200:
        return True
    # a generated-file diff (a/... -> b/... plus an @@ hunk header)
    if re.search(r"\ba//\S+\s*→\s*b//\S+\b", text) and re.search(r"^@@[^@]*@@$", text, flags=re.M):
        return True
    return False


@router.post("/telegram")
async def telegram(request: Request, background_tasks: BackgroundTasks):
    # Get client info for audit logging
    client_ip = request.headers.get("x-forwarded-for", "").split(",")[0].strip() or \
                (request.client.host if request.client else "unknown")
    request_id = getattr(request.state, "request_id", None)

    # Verify webhook secret if configured
    if settings.telegram_webhook_secret:
        secret_header = request.headers.get("X-Telegram-Bot-Api-Secret-Token", "")
        if secret_header != settings.telegram_webhook_secret:
            await audit_logger.log_event(
                event_type=AuditLogger.EventType.WEBHOOK_SIGNATURE_FAILED,
                severity=AuditLogger.Severity.WARNING,
                ip_address=client_ip,
                request_id=request_id,
                details={"platform": "telegram", "reason": "invalid_webhook_secret"},
            )
            logger.warning("telegram_webhook_invalid_secret", ip_address=client_ip)
            raise HTTPException(403, "Invalid webhook secret")

    body = await request.json()
    chat_id = str(body.get("message", {}).get("chat", {}).get("id", ""))
    text_msg = body.get("message", {}).get("text", "")
    chat_type = body.get("message", {}).get("chat", {}).get("type", "private")
    sender_id = str(body.get("message", {}).get("from", {}).get("id", ""))
    if not chat_id or not text_msg:
        return {"status": "ok"}

    # ── SECURITY: lock down who can trigger an agent ──
    # Only accept private 1:1 messages. Groups/channels (a shared "beprepared channel")
    # can be posted to by anyone, so they must NEVER route to a profile.
    if chat_type not in ("private", "supergroup_private"):
        logger.warning("telegram_chat_blocked", chat_type=chat_type, chat_id=chat_id)
        await audit_logger.log_event(
            event_type=AuditLogger.EventType.WEBHOOK_MESSAGE_RECEIVED,
            severity=AuditLogger.Severity.WARNING,
            ip_address=client_ip,
            request_id=request_id,
            details={"platform": "telegram", "reason": "non_private_chat_blocked", "chat_type": chat_type, "chat_id": chat_id},
        )
        return {"status": "ignored"}
    # In a private chat the sender id equals the chat id. Reject any update where a
    # different account is pretending to be the chat owner.
    if sender_id and sender_id != chat_id:
        logger.warning("telegram_sender_blocked", chat_id=chat_id, sender_id=sender_id)
        await audit_logger.log_event(
            event_type=AuditLogger.EventType.WEBHOOK_MESSAGE_RECEIVED,
            severity=AuditLogger.Severity.WARNING,
            ip_address=client_ip,
            request_id=request_id,
            details={"platform": "telegram", "reason": "sender_mismatch_blocked", "chat_id": chat_id, "sender_id": sender_id},
        )
        return {"status": "ignored"}

    async with async_session_factory() as db:
        r = await db.execute(
            text("SELECT id,is_active,profile_path,runtime FROM user_profiles WHERE phone_number=:c"),
            {"c": chat_id},
        )
        u = r.fetchone()

        # ── Handle /start commands for ALL users (existing or not) ──
        if text_msg.startswith("/start "):
            code = text_msg.split(" ", 1)[1].strip()

            # Handle telegram-link codes (link existing user to Telegram)
            if code.startswith("link_"):
                r = await db.execute(
                    text("SELECT claimed_by, agent_name FROM invite_links WHERE code=:c AND claimed_by IS NOT NULL"),
                    {"c": code},
                )
                link = r.fetchone()
                if not link:
                    r2 = await db.execute(text("SELECT claimed_by FROM invite_links WHERE code=:c"), {"c": code})
                    if r2.fetchone():
                        await send_tg(chat_id, "✅ You're already connected! Send me any message.")
                    else:
                        await send_tg(chat_id, "❌ This link is invalid or expired.")
                    return {"status": "ok"}
                uid = str(link[0])
                name = link[1] or "Agent"
                await db.execute(text("UPDATE user_profiles SET phone_number=:c, platform='telegram' WHERE id::text=:uid"), {"c": chat_id, "uid": uid})
                await db.execute(text("DELETE FROM invite_links WHERE code=:c"), {"c": code})
                await db.commit()
                try:
                    from pathlib import Path
                    if not (Path("/opt/hermes/profiles") / uid / "config.yaml").exists():
                        profile = init_user_profile(user_id=uid, agent_name=name, plan="pro")
                        await db.execute(text("UPDATE user_profiles SET profile_path=:pp WHERE id::text=:uid"), {"pp": profile["profile_dir"], "uid": uid})
                except:
                    pass
                await send_tg(chat_id,
                    f"✅ *Connected!* Your Telegram is now linked to **{name}**.\n\n"
                    f"Try saying:\n"
                    f"📝 *\"Save a note about...\"*\n"
                    f"🌐 *\"Search for...\"*\n"
                    f"📖 *\"What's in my vault?\"*")
                return {"status": "linked"}

            # Handle invite codes (new users only)
            if not u:
                r = await db.execute(
                    text("SELECT id,label,agent_name,plan,trial_days,is_vip FROM invite_links WHERE code=:c AND claimed_by IS NULL"),
                    {"c": code},
                )
                inv = r.fetchone()
                if not inv:
                    await send_tg(chat_id, "❌ This invite link is invalid or has already been used.")
                    return {"status": "invalid"}
                agent_name = inv[2] or "My Assistant"
                plan = inv[3] or "pro"
                is_vip = inv[5]
                trial_ends = datetime.utcnow() + timedelta(days=inv[4]) if inv[4] else None
                r2 = await db.execute(
                    text("""INSERT INTO user_profiles (phone_number, agent_name, plan, is_vip, trial_ends_at, primary_model, backup_model, platform)
                           VALUES (:p,:a,:pl,:v,:te,:m1,:m2,'telegram') RETURNING id"""),
                    {"p": chat_id, "a": agent_name, "pl": plan, "v": is_vip, "te": trial_ends,
                     "m1": settings.default_primary_model, "m2": settings.default_backup_model},
                )
                uid = str(r2.fetchone()[0])
                try:
                    profile = init_user_profile(user_id=uid, agent_name=agent_name, plan=plan, is_vip=is_vip)
                    await db.execute(
                        text("UPDATE user_profiles SET profile_path=:pp WHERE id::text=:uid"),
                        {"pp": profile["profile_dir"], "uid": uid},
                    )
                except Exception as e:
                    pass
                await db.execute(
                    text("UPDATE invite_links SET claimed_by=:u, claimed_at=NOW() WHERE code=:c"),
                    {"u": uid, "c": code},
                )
                await db.commit()
                await send_tg(chat_id, WELCOME)
                return {"status": "activated"}
            else:
                await send_tg(chat_id, "👋 I don't know you yet! You need an invite link to use this bot.")
                return {"status": "ignored"}

        if not u or not u[1]:
            await send_tg(chat_id, "⏳ Your account is inactive. Contact your admin.")
            return {"status": "inactive"}

        # ── Process message with typing indicator ──
        stop_typing = asyncio.Event()
        typing_task = asyncio.create_task(typing_indicator(chat_id, stop_typing))

        # Handle document uploads (PDFs, docs, etc.) to knowledge base
        doc = body.get("message", {}).get("document", None)
        if doc:
            try:
                file_id = doc.get("file_id", "")
                file_name_raw = doc.get("file_name", "document")
                # SECURITY FIX: Sanitize filename to prevent path traversal attacks
                # Use os.path.basename to strip any directory components
                file_name = os.path.basename(file_name_raw)
                # Additional validation: reject suspicious characters
                if not file_name or '..' in file_name or '/' in file_name or '\\' in file_name:
                    file_name = f"document_{int(__import__('time').time())}"
                # Ensure the file has a reasonable extension
                if not any(file_name.lower().endswith(ext) for ext in ['.pdf', '.doc', '.docx', '.txt', '.md', '.csv', '.xlsx', '.pptx', '.png', '.jpg', '.jpeg']):
                    file_name = file_name + '.dat'
                async with httpx.AsyncClient() as c:
                    fr = await c.get(f"https://api.telegram.org/bot{settings.telegram_bot_token}/getFile?file_id={file_id}")
                    fp = fr.json().get("result", {}).get("file_path", "")
                    if fp:
                        dl = await c.get(f"https://api.telegram.org/bot{settings.telegram_bot_token}/{fp}")
                        uid = str(u[0])
                        kb_dir = Path("/opt/hermes/obsidian") / uid / "Knowledge"
                        kb_dir.mkdir(parents=True, exist_ok=True)
                        # Double-check the final path is within kb_dir (defense in depth)
                        save_path = (kb_dir / file_name).resolve()
                        if not str(save_path).startswith(str(kb_dir.resolve())):
                            raise ValueError("Invalid file path")
                        save_path.write_bytes(dl.content)

                        # Log document upload
                        await audit_logger.log_event(
                            event_type=AuditLogger.EventType.WEBHOOK_DOCUMENT_UPLOAD,
                            severity=AuditLogger.Severity.INFO,
                            user_id=uid,
                            ip_address=client_ip,
                            request_id=request_id,
                            details={"platform": "telegram", "file_name": file_name, "file_size": len(dl.content)},
                        )
                        logger.info("telegram_document_uploaded", user_id=uid, file_name=file_name, file_size=len(dl.content))

                        await send_tg(chat_id, f"📎 Saved **{file_name}** to your knowledge base.")
                        stop_typing.set()
                        typing_task.cancel()
                        return {"status": "document_saved"}
            except Exception as e:
                logger.error("telegram_document_upload_failed", user_id=str(u[0]), error=str(e), exc_info=True)
                await send_tg(chat_id, "⚠️ Couldn't save the document. Please try again.")
                stop_typing.set()
                typing_task.cancel()
                return {"status": "doc_error"}

        # Handle voice messages — transcribe and process as text
        voice = body.get("message", {}).get("voice", None)
        if voice:
            try:
                file_id = voice.get("file_id", "")
                async with httpx.AsyncClient() as c:
                    fr = await c.get(f"https://api.telegram.org/bot{settings.telegram_bot_token}/getFile?file_id={file_id}")
                    fp = fr.json().get("result", {}).get("file_path", "")
                    if fp:
                        dl = await c.get(f"https://api.telegram.org/bot{settings.telegram_bot_token}/{fp}")
                        # Save OGG, convert to WAV, transcribe
                        ogg_path = Path(f"/tmp/voice_{chat_id}.ogg")
                        wav_path = Path(f"/tmp/voice_{chat_id}.wav")
                        ogg_path.write_bytes(dl.content)
                        # Convert to WAV using pydub
                        from pydub import AudioSegment
                        audio = AudioSegment.from_ogg(str(ogg_path))
                        audio.export(str(wav_path), format="wav")
                        # Transcribe using SpeechRecognition
                        import speech_recognition as sr
                        recognizer = sr.Recognizer()
                        with sr.AudioFile(str(wav_path)) as source:
                            audio_data = recognizer.record(source)
                            text_msg = recognizer.recognize_google(audio_data)
                        # Cleanup temp files
                        ogg_path.unlink(missing_ok=True)
                        wav_path.unlink(missing_ok=True)
                        await send_tg(chat_id, f"🎤 *You said:* {text_msg}")
                # Fall through to regular agent processing with transcribed text_msg
            except Exception as e:
                await send_tg(chat_id, "⚠️ Couldn't transcribe your voice message. Please try again or type instead.")
                stop_typing.set()
                typing_task.cancel()
                return {"status": "voice_error"}

        # ── Agent personality (SOUL) commands ──
        _low = (text_msg or "").strip()
        if _low.startswith("/personality") or _low.startswith("/soul"):
            _rest = _low.split(None, 1)[1] if " " in _low else ""
            if _rest:
                async with async_session_factory() as db:
                    await db.execute(text("UPDATE user_profiles SET personality=:p WHERE id::text=:u"), {"p": _rest, "u": str(u[0])})
                    await db.commit()
                await send_tg(chat_id, "✅ **Personality updated!** I'll follow these instructions from now on.\n\n" + _rest[:400])
            else:
                async with async_session_factory() as db:
                    _r = await db.execute(text("SELECT personality FROM user_profiles WHERE id::text=:u"), {"u": str(u[0])})
                    _row = _r.fetchone()
                if _row and _row[0]:
                    await send_tg(chat_id, "🧠 **Your agent personality:**\n\n" + _row[0][:1000])
                else:
                    await send_tg(chat_id, "🧠 Your agent uses the **default personality**. Send `/personality` followed by your instructions to customize me — e.g. `/personality Always keep replies under 100 words and greet me by name.`")
            stop_typing.set()
            typing_task.cancel()
            return {"status": "personality"}

        # PERFORMANCE FIX: Always use background task to avoid blocking webhook response
        # This prevents Telegram from timing out and retrying, and gives users faster acknowledgment
        stop_typing.set()
        typing_task.cancel()

        # Determine if this is a Hermes runtime user
        is_hermes_runtime = len(u) > 3 and u[3] == "hermes"
        profile_dir = str(u[2]) if u[2] else None

        background_tasks.add_task(
            _process_message_async,
            str(u[0]),
            chat_id,
            text_msg,
            profile_dir,
            is_hermes_runtime,
            client_ip,
            request_id,
        )

        await audit_logger.log_event(
            event_type=AuditLogger.EventType.WEBHOOK_MESSAGE_RECEIVED,
            severity=AuditLogger.Severity.INFO,
            user_id=str(u[0]),
            ip_address=client_ip,
            details={"note": f"[async] {text_msg[:120]}", "runtime": "hermes" if is_hermes_runtime else "default"},
        )

        # Log activity asynchronously too
        background_tasks.add_task(
            _log_activity,
            str(u[0]),
            text_msg,
            request_id,
            client_ip,
        )

        return {"status": "ok", "async": True}


_user_locks: dict = {}


def _get_user_lock(user_id: str):
    """Per-user asyncio lock. Each user gets their own lock so concurrent messages
    for the SAME user can't race on that user's (per-profile) Hermes session.
    Distinct users have distinct locks + distinct session dirs -> no cross-user.
    """
    lk = _user_locks.get(user_id)
    if lk is None:
        lk = asyncio.Lock()
        _user_locks[user_id] = lk
    return lk


async def _read_stream(stream, sink, prog):
    while True:
        chunk = await stream.read(2048)
        if not chunk:
            break
        sink.append(chunk)
        prog["last"] = time.monotonic()
        prog["bytes"] += len(chunk)


async def _run_hermes(args: list, timeout: int, stall_s: int = 55):
    """Spawn hermes once. Returns (stdout, stderr, returncode).

    Progress-aware: while hermes produces output (each API call prints a line) it
    runs; if it goes SILENT for > `stall_s` (a model call hung), we kill it and
    return -1 so the caller can fail fast instead of stranding the user ~3min.
    """
    proc = await asyncio.create_subprocess_exec(
        *args,
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
        env={**os.environ, "HERMES_HOME": "/opt/hermes/hermes"},
        cwd="/opt/hermes/hermes",
    )
    out, err = [], []
    prog = {"last": time.monotonic(), "bytes": 0, "deadline": time.monotonic() + timeout}
    rt = asyncio.create_task(_read_stream(proc.stdout, out, prog))
    re = asyncio.create_task(_read_stream(proc.stderr, err, prog))
    try:
        rc = None
        while rc is None:
            try:
                rc = await asyncio.wait_for(proc.wait(), timeout=3)
            except asyncio.TimeoutError:
                now = time.monotonic()
                if now > prog["deadline"]:
                    break  # total cap
                if now - prog["last"] > stall_s:
                    try:
                        proc.kill()
                    except Exception:
                        pass
                    rc = await proc.wait()
                    break
        rt.cancel()
        re.cancel()
        return b"".join(out), b"".join(err), rc if rc is not None else -1
    except Exception:
        try:
            proc.kill()
        except Exception:
            pass
        rt.cancel()
        re.cancel()
        return None, b"", -1


async def run_hermes_runtime(user_id: str, message: str, timeout: int = 360) -> str | None:
    """Invoke the container's Hermes runtime headlessly for a profile.

    Returns the assistant reply, or None if Hermes is unavailable/fails so the
    caller can gracefully fall back to the platform's agent_manager. Safe by default.
    """
    import asyncio, shutil
    herm = shutil.which("hermes")
    if not herm:
        return None
    try:
        _t = time.perf_counter()
        lock = _get_user_lock(user_id)
        async with lock:
            # Resume the user's OWN latest session (scoped by -p <uuid> to their
            # profile dir) so the system-prompt prefix is KV-cache reused and the
            # big first prefill becomes a cache hit. Per-user lock prevents two
            # concurrent messages for the same user racing on their session.
            out, _err, rc = await _run_hermes(
                [herm, "-p", user_id, "chat", "-r", "latest", "-q", message,
                 "-Q", "--reasoning", "none"], timeout)
            if rc == -1:
                # A model call stalled (watchdog killed it). Don't burn another
                # ~55s on a fresh retry - return None so the caller tells the
                # user promptly instead of stranding them for minutes.
                logger.info("hermes_stall", user_id=user_id, stage="resume")
                return None
            if out is None or rc != 0:
                # No session yet (first turn) or a transient error -> fresh session.
                out, _err, rc = await _run_hermes(
                    [herm, "-p", user_id, "chat", "-q", message, "-Q", "--reasoning", "none"], timeout)
                if rc == -1:
                    logger.info("hermes_stall", user_id=user_id, stage="fresh")
                    return None
        _calls = "; ".join(l.strip() for l in (_err or b"").decode("utf-8", "replace").splitlines() if "latency=" in l)
        logger.info("hermes_turn", user_id=user_id, elapsed=f"{time.perf_counter()-_t:.1f}s", calls=_calls or "n/a", out_len=len(out or b""))
        text = (out or b"").decode("utf-8", "replace")
        # -Q emits headers/session-summary lines followed by the actual response
        # (plain text OR a deck-diff the user wants as .pptx). Drop the CLI chrome
        # lines but KEEP the full response so deck JSON isn't truncated away.
        resp = "\n".join(
            ln for ln in text.splitlines()
            if not ln.strip().startswith((
                "session_id:", "Session:", "Title:", "Duration:", "Messages:",
                "Query:", "Initializingagent", "Initializing agent",
                "┌", "└", "┐", "┘", "│", "╭", "╮", "╰", "╯", "─", "═",
                "Resume this session", "hermes --resume", "hermes -c",
            ))
        ).strip() or None
        return resp
    except Exception:
        return None


_FAST_IDENT = {
    "events": re.compile(r"\b(schedule|calendar|events?|appointments?|upcoming)\b", re.I),
    "reminders": re.compile(r"\b(remind(er|ers)?|to[- ]do|todos?)\b", re.I),
    "notes": re.compile(r"\bnotes?\b", re.I),
    "projects": re.compile(r"\bprojects?\b", re.I),
}
_FAST_WRITE = re.compile(
    r"\b(create|add|new|set|delete|remove|update|edit|make|write|"
    r"book|cancel|rename|reschedule|postpone|schedule a|remind me to|remind me of)\b", re.I)


def _detect_fast_kind(low: str):
    for k, pat in _FAST_IDENT.items():
        if pat.search(low):
            return k
    return None


def _fmt_dt(dt):
    if dt is None:
        return "sometime"
    try:
        return dt.astimezone(ZoneInfo("America/New_York")).strftime("%a %b %d  %I:%M %p")
    except Exception:
        return str(dt)[:16]


async def _fast_query(kind: str, user_id: str):
    async with async_session_factory() as db:
        if kind == "events":
            r = await db.execute(text(
                "SELECT title, event_start FROM scheduled_events WHERE user_id::text=:u AND event_start >= NOW() ORDER BY event_start LIMIT 8"),
                {"u": user_id})
            return [f"• {t} — {_fmt_dt(s)}" for t, s in r.fetchall()]
        if kind == "reminders":
            r = await db.execute(text(
                "SELECT title, remind_at FROM reminders WHERE user_id::text=:u AND done=false ORDER BY remind_at LIMIT 8"),
                {"u": user_id})
            return [f"• {t} — {_fmt_dt(s)}" for t, s in r.fetchall()]
        if kind == "notes":
            r = await db.execute(text(
                "SELECT title, content FROM notes WHERE user_id::text=:u ORDER BY updated_at DESC LIMIT 6"),
                {"u": user_id})
            out = []
            for t, c in r.fetchall():
                c = (c or "").strip().replace("\n", " ")[:80]
                out.append(f"• {t}" + (f": {c}…" if c else ""))
            return out
        if kind == "projects":
            r = await db.execute(text(
                "SELECT title FROM projects WHERE user_id::text=:u ORDER BY updated_at DESC LIMIT 8"),
                {"u": user_id})
            return [f"• {t}" for (t,) in r.fetchall()]
    return []


async def _fast_count(kind: str, user_id: str) -> int:
    tbl = {"events": "scheduled_events", "reminders": "reminders",
           "notes": "notes", "projects": "projects"}[kind]
    async with async_session_factory() as db:
        r = await db.execute(text(f"SELECT COUNT(*) FROM {tbl} WHERE user_id::text=:u"), {"u": user_id})
        return r.scalar() or 0


async def try_fast_path(user_id: str, message: str):
    """Answer direct DB-lookup questions instantly (no agent loop, no model calls).
    Returns a reply string, or None to fall through to the Hermes agent."""
    m = (message or "").strip()
    low = m.lower()
    if not m or len(m) > 200:
        return None
    if _FAST_WRITE.search(low):
        return None                                     # actions/writes -> agent
    if re.search(r"\b(how many)\b|\bcount of\b", low):
        k = _detect_fast_kind(low)
        if k:
            return f"You have {await _fast_count(k, user_id)} {k}."
        return None
    k = _detect_fast_kind(low)
    if not k:
        return None
    if not (re.search(r"\b(what|list|show|tell|see|my|any|current|upcoming|all|how)\b", low)
            or k in ("notes", "projects")):
        return None
    items = await _fast_query(k, user_id)
    head = {"events": "upcoming events", "reminders": "active reminders",
            "notes": "notes", "projects": "projects"}[k]
    if not items:
        return f"You have no {head} right now."
    reply = f"Here are your {head}:\n" + "\n".join(items)
    if len(items) == 8:
        reply += f"\n…showing the {len(items)} most recent."
    return reply


async def _run_hermes_async(user_id: str, chat_id: str, message: str, client_ip: str, request_id: str) -> None:
    """Background delivery for the async Hermes turn.

    Runs the per-user Hermes runtime (bridge-enabled), falls back to the platform
    agent_manager if Hermes is unavailable, then posts the reply to Telegram.
    """
    # Show typing immediately (typing_indicator sends on its first iteration) so the
    # user sees the bot working.
    typing_stop = asyncio.Event()
    typing_task = asyncio.create_task(typing_indicator(chat_id, typing_stop))
    try:
        # Fast-path: answer direct DB-lookup questions instantly (no agent/model).
        quick = await try_fast_path(user_id, message)
        if quick:
            typing_stop.set()
            typing_task.cancel()
            await send_tg(chat_id, quick)
            await audit_logger.log_event(
                event_type=AuditLogger.EventType.WEBHOOK_MESSAGE_RECEIVED,
                severity=AuditLogger.Severity.INFO,
                user_id=user_id,
                ip_address=client_ip,
                details={"note": f"[fast-path] {quick[:140]}"},
            )
            return
        async def _slow_ack():
            # If a turn drags past ~8s, reassure the user it's processing and not
            # ignored (a silent 2-3 min wait reads as "you're ignoring me").
            await asyncio.sleep(8)
            try:
                await send_tg(chat_id, "Just so you know — I'm still working on this; it's taking a bit longer. I'll reply shortly.")
            except Exception:
                pass
        ack_task = asyncio.create_task(_slow_ack())
        try:
            resp = await run_hermes_runtime(user_id, message)
        finally:
            ack_task.cancel()
        if resp is None:
            # Hermes failed/stalled/timed out. Be honest with the user instead of
            # silently falling back to a slower, weaker legacy path (or a blank).
            resp = ("I hit a temporary hiccup and couldn't finish that request (it's on my "
                    "end, not yours). Please try again in a moment — sorry!")
        if resp:
            await deliver_reply(chat_id, resp)
        else:
            # Never leave the user hanging on an empty response.
            await send_tg(chat_id, "Hmm, I couldn't turn anything up for that — can you rephrase?")
        await audit_logger.log_event(
            event_type=AuditLogger.EventType.WEBHOOK_MESSAGE_RECEIVED,
            severity=AuditLogger.Severity.INFO,
            user_id=user_id,
            ip_address=client_ip,
            details={"note": f"[hermes reply] {str(resp)[:160]}"},
        )
    except Exception:
        logger.exception("hermes_async_failed")
        try:
            await send_tg(chat_id, "⚠️ Something went wrong while I was processing your message. Please try again.")
        except Exception:
            pass
    finally:
        typing_stop.set()
        typing_task.cancel()


async def _process_message_async(
    user_id: str,
    chat_id: str,
    message: str,
    profile_dir: str | None,
    is_hermes_runtime: bool,
    client_ip: str,
    request_id: str,
) -> None:
    """Unified background message processor for all users.

    Routes to Hermes runtime or platform agent_manager based on user config,
    with fallback handling. Sends typing indicator while processing.
    """
    typing_stop = asyncio.Event()
    typing_task = asyncio.create_task(typing_indicator(chat_id, typing_stop))
    try:
        # Fast-path: answer direct DB-lookup questions instantly (0 model calls).
        quick = await try_fast_path(user_id, message)
        if quick:
            typing_stop.set()
            typing_task.cancel()
            await send_tg(chat_id, quick)
            await audit_logger.log_event(
                event_type=AuditLogger.EventType.WEBHOOK_MESSAGE_RECEIVED,
                severity=AuditLogger.Severity.INFO,
                user_id=user_id,
                ip_address=client_ip,
                details={"note": f"[fast-path] {quick[:140]}"},
            )
            return

        resp = None
        if is_hermes_runtime:
            # Try Hermes runtime first for hermes users
            resp = await run_hermes_runtime(user_id, message)

        # Fallback to platform agent_manager
        if resp is None:
            resp = await hermes_profile_chat_with_fallback(
                user_id=user_id,
                message=message,
                profile_dir=profile_dir,
            )

        if resp:
            await deliver_reply(chat_id, resp)

        logger.info("message_processed", user_id=user_id, message_length=len(message), runtime="hermes" if is_hermes_runtime else "default")

    except Exception as e:
        logger.exception("message_processing_failed", user_id=user_id, error=str(e))
        try:
            await send_tg(chat_id, "⚠️ Something went wrong. Please try again in a moment.")
        except Exception:
            pass
    finally:
        typing_stop.set()
        typing_task.cancel()


async def _log_activity(user_id: str, message: str, request_id: str, client_ip: str) -> None:
    """Background task to log activity without blocking webhook response."""
    try:
        async with async_session_factory() as db:
            # SECURITY FIX: Use json.dumps to prevent JSON injection
            await db.execute(
                text("INSERT INTO activity_logs (user_id,action,details,request_id,ip_address) VALUES (:uid,'message',:det,:rid,:ip)"),
                {"uid": user_id, "det": json.dumps({"platform": "telegram", "tokens": len(message) // 4}), "rid": request_id, "ip": client_ip},
            )
            await db.commit()
    except Exception as e:
        logger.warning("activity_log_failed", user_id=user_id, error=str(e))